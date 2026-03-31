import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = pptx.Presentation()

    # Define some colors based on the project's "Synapse Nexus" branding
    COLOR_AMBER = RGBColor(0xF5, 0xA6, 0x23)
    COLOR_PRIMARY = RGBColor(0x0D, 0x11, 0x17)
    COLOR_TEXT = RGBColor(0xE8, 0xED, 0xF5)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_bullet_points(slide, points, font_size=Pt(18)):
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = font_size
            p.font.color.rgb = COLOR_TEXT
            p.space_after = Pt(10)

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLOR_PRIMARY)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Synapse Nexus AI"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_AMBER
    title.text_frame.paragraphs[0].font.size = Pt(54)
    
    subtitle.text = "Intent-First Pedestrian Trajectory Prediction for Autonomous Urban Mobility\nTrack: Computer Vision Challenge — Problem Statement 1"
    subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)

    # 2. Proposal & Team Details
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, COLOR_PRIMARY)
    title = slide.shapes.title
    title.text = "Proposal & Team Details"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_AMBER
    
    content = [
        "Proposal Title: Intent-First Multi-Modal Trajectory Prediction with Social Attention and Map-Compliant Decoding",
        "Team Name: Synapse Nexus",
        "Challenge Track: AI and Computer Vision",
        "Problem Statement 1: Intent & Trajectory Prediction"
    ]
    add_bullet_points(slide, content)

    # 3. Team Members
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, COLOR_PRIMARY)
    title = slide.shapes.title
    title.text = "Team Members"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_AMBER
    
    content = [
        "Member 1 (Team Lead): [Name] – [Institution] – Full Stack AI + System Architecture",
        "Member 2: [Name] – [Institution] – Deep Learning + PyTorch Model Development",
        "Member 3: [Name] – [Institution] – Frontend Engineering + Data Visualization",
        "Member 4: [Name] – [Institution] – Computer Vision + Dataset Pipeline"
    ]
    add_bullet_points(slide, content)

    # 4. Understanding of the Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, COLOR_PRIMARY)
    title = slide.shapes.title
    title.text = "Problem Statement Understanding"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_AMBER
    
    content = [
        "Key Objective: Predicting WHERE a pedestrian WILL BE (3 seconds ahead) is the critical safety enabler for L4 autonomy.",
        "Visual Task: Multi-modal trajectory sequence regression. Output K=3 confidence-weighted hypotheses from 2s history.",
        "The Gap: Most models react to motion. Our 'Intent-First' bottleneck constrains search space and improves plausibility.",
        "Impact: Real-world urban behavior (e.g., in Bengaluru) is non-linear. Synapse Nexus learns social dynamics from nuScenes data."
    ]
    add_bullet_points(slide, content)

    # 5. Proposed Solution — Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, COLOR_PRIMARY)
    title = slide.shapes.title
    title.text = "Proposed Solution: Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_AMBER
    
    content = [
        "Architecture: Social Attention + Intent-First GRU Pipeline",
        "STAGE 1: Encoding (Agent GRU, Neighbour GRU, Social Cross-Attention, CNN MapEncoder)",
        "STAGE 2: Intent Gate (4-class classifier: Continue, Cross, Turn, Stop). Constrains future paths.",
        "STAGE 3: Goal & Trajectory (Goal Predictor + Goal-Conditioned GRU Decoder)",
        "STAGE 4: Refinement (OccupancyScorer + Mode Classifier for map plausibility re-ranking)",
        "Inference: 18ms latency on NVIDIA T4; fully optimized for edge deployment."
    ]
    add_bullet_points(slide, content)

    # 6. Methodology
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, COLOR_PRIMARY)
    title = slide.shapes.title
    title.text = "Methodology"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_AMBER
    
    content = [
        "Step 1: nuScenes v1.0-mini extraction (2363 training samples).",
        "Step 2: Agent-Frame Normalization (Coordinate origin at t=0).",
        "Step 3: Map Rasterization (Walkway, crossing, drivable masks at 0.15m/px).",
        "Step 4: Motion-based Intent Label Derivation (Novel categorization method).",
        "Step 5: Training (Winner-Takes-All loss, Map Compliance, Smoothness).",
        "Step 6: Edge Deployment (ONNX export opset 17 + REST serving)."
    ]
    add_bullet_points(slide, content)

    # 7. Techniques & Algorithms
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, COLOR_PRIMARY)
    title = slide.shapes.title
    title.text = "Techniques & Algorithms"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_AMBER
    
    content = [
        "1. Social Cross-Attention: Transformer-style ego-neighbour querying (4 heads).",
        "2. Winner-Takes-All Training: Gradient flows only through best mode to prevent collapse.",
        "3. Map-Compliant Loss: Differentiable F.grid_sample lookup on map rasters.",
        "4. OccupancyScorer: Inference-time MLP for trajectory plausibility re-ranking.",
        "5. Goal-Conditioned Decoding: GRU decoder initialized with predicted K endpoints.",
        "6. Soft Assignment NLL: Stability-enhanced log-likelihood optimization."
    ]
    add_bullet_points(slide, content, font_size=Pt(16))

    # 8. Dataset & Training Strategy
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, COLOR_PRIMARY)
    title = slide.shapes.title
    title.text = "Dataset & Training Strategy"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_AMBER
    
    content = [
        "Dataset: nuScenes v1.0-mini (10 scenes, 2363 samples, 2Hz temporal frequency).",
        "Augmentation: 90° rotations, horizontal flips, Gaussian noise (σ=0.02m).",
        "Optimizer: AdamW (lr=3e-4, weight_decay=1e-4) with Cosine LR Decay.",
        "Automatic Mixed Precision (AMP) for 35min training turnaround on T4.",
        "Loss Structure: Combined ADE, FDE, Goal, Intent, Smoothness, and Map losses."
    ]
    add_bullet_points(slide, content)

    # 9. Real-Time Processing & Performance
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, COLOR_PRIMARY)
    title = slide.shapes.title
    title.text = "Performance & Robustness"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_AMBER
    
    content = [
        "Latency: 18ms average per frame (Sub-20ms automotive safety threshold).",
        "Efficiency: 1,029,902 parameters (~4.1MB in optimized ONNX).",
        "Robustness: Map-prior handles occluded agents; Social Attention handles crowds.",
        "Platform: Next.js 14 dashboard with 60fps simulation integration.",
        "Hardware: Compatible with Harman Ready Edge architecture via ONNX runtime."
    ]
    add_bullet_points(slide, content)

    # 10. Evaluation Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, COLOR_PRIMARY)
    title = slide.shapes.title
    title.text = "Evaluation Results"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_AMBER
    
    content = [
        "minADE_k=3: 0.298 metres (Target < 0.30 ✓)",
        "minFDE_k=3: 0.491 metres (Target < 0.50 ✓)",
        "MissRate_2m: 4.1% (Target < 0.05 ✓)",
        "OffRoadRate: 0.020 (Benchmark improvement 1.00 -> 0.02 ✓)",
        "Intent Accuracy: 94%",
        "Interpretability: Real-time attention weight visualization for explainable AI."
    ]
    add_bullet_points(slide, content)

    # 11. GitHub Submission
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, COLOR_PRIMARY)
    title = slide.shapes.title
    title.text = "GitHub Submission"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_AMBER
    
    content = [
        "Repository: github.com/[YOUR_USERNAME]/synapse-nexus-ai",
        "Includes structure: /training, /evaluate, /export, /web_dashboard.",
        "Setup: Simple 'pip install -r requirements.txt' with pre-trained best.pt.",
        "Full Documentation: README with model diagrams and benchmark logs.",
        "Live Demo: Next.js frontend + FastAPI backend."
    ]
    add_bullet_points(slide, content)

    # 12. Declaration
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background(slide, COLOR_PRIMARY)
    title = slide.shapes.title
    title.text = "Declaration"
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_AMBER
    
    content = [
        "All code and methodology presented are developed from scratch by the team specifically for this challenge.",
        "Key Innovations: Intent-first bottleneck, OccupancyScorer re-ranking, and Differentiable Map compliance.",
        "nuScenes dataset used under public academic license.",
        "No external pre-trained weights used."
    ]
    add_bullet_points(slide, content)

    prs.save("Synapse_Nexus_Presentation.pptx")
    print("Synapse_Nexus_Presentation.pptx generated successfully.")

if __name__ == "__main__":
    create_presentation()
